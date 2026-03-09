#!/usr/bin/env python3
"""Convert slide PNG images in an output directory into a single PPTX file."""

import argparse
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from PIL import Image


# 2K slide dimensions: 2752x1536 px at 96 dpi -> inches -> EMU
# Standard widescreen: 13.33 x 7.5 inches (33.867 x 19.05 cm)
SLIDE_WIDTH_EMU  = 12192000   # 13.33 in * 914400
SLIDE_HEIGHT_EMU = 6858000    # 7.5  in * 914400


def images_to_pptx(images_dir: str, output_pptx: str) -> None:
    images_dir = Path(images_dir)
    image_files = sorted(images_dir.glob("slide-*.png"))

    if not image_files:
        print(f"No slide images found in {images_dir}")
        sys.exit(1)

    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    for img_path in image_files:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"  Added: {img_path.name}")

    prs.save(output_pptx)
    print(f"\nPPTX saved: {output_pptx}")


def main():
    parser = argparse.ArgumentParser(description="Convert slide PNGs to PPTX")
    parser.add_argument("--images-dir", required=True, help="Directory containing slide-XX.png files")
    parser.add_argument("--output",     required=True, help="Output .pptx file path")
    args = parser.parse_args()

    images_to_pptx(args.images_dir, args.output)


if __name__ == "__main__":
    main()
